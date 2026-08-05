# -*- coding: utf-8 -*-
"""Genera los 8 activos ficticios MCP-365-S2-v2 · Proyecto Horizonte.

Salida: planillas/01_…08_… y ZIP MCP365_S2_Kit_Proyecto_Horizonte.zip
Fines académicos · Datos ficticios.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    HRFlowable,
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "planillas"

DARK = "211D40"
PURPLE = "4B3FAA"
YELLOW = "FFEC00"
LIGHT = "F4F3F9"
MUTED = "5A5A72"
WHITE = "FFFFFF"
ACADEMIC = "Fines académicos · Datos ficticios"

FILENAMES = [
    "01_Correo_Solicitud_Proyecto_Horizonte.pdf",
    "02_Alcance_Proyecto_Horizonte.docx",
    "03_Presupuesto_y_Cronograma_Horizonte.xlsx",
    "04_Transcripcion_Reunion_Horizonte.docx",
    "05_Registro_Inicial_Riesgos_Horizonte.xlsx",
    "06_Comentarios_Interesados_Horizonte.docx",
    "07_Plantilla_Comite_Horizonte.pptx",
    "08_Guia_Validacion_Resultados.pdf",
]


# ---------------------------------------------------------------------------
# Helpers docx
# ---------------------------------------------------------------------------
def shade_cell(cell, hex_color: str) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for child in list(tcPr):
        if child.tag == qn("w:shd"):
            tcPr.remove(child)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tcPr.append(shd)


def set_run(run, *, size=11, bold=False, color=None, font="Calibri", italic=False) -> None:
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    if color:
        run.font.color.rgb = RGBColor.from_string(color)


def add_para(
    doc,
    text,
    *,
    size=11,
    bold=False,
    color=None,
    space_after=6,
    space_before=0,
    align=None,
    italic=False,
):
    p = doc.add_paragraph()
    if align:
        p.alignment = align
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.space_before = Pt(space_before)
    run = p.add_run(text)
    set_run(run, size=size, bold=bold, color=color or DARK, italic=italic)
    return p


def add_heading_styled(doc, text: str) -> None:
    add_para(doc, text, size=13, bold=True, color=PURPLE, space_before=12, space_after=6)


def add_bullet(doc, text: str, *, size=10) -> None:
    p = doc.add_paragraph(style="List Bullet")
    p.clear()
    run = p.add_run(text)
    set_run(run, size=size, color=DARK)
    p.paragraph_format.space_after = Pt(3)


def add_banner(doc, title: str, subtitle: str = "") -> None:
    table = doc.add_table(rows=1, cols=1)
    table.autofit = True
    cell = table.rows[0].cells[0]
    shade_cell(cell, DARK)
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(title)
    set_run(run, size=14, bold=True, color=YELLOW)
    if subtitle:
        p2 = cell.add_paragraph()
        run2 = p2.add_run(subtitle)
        set_run(run2, size=9, color="D5D2E6")
    add_para(doc, "", space_after=8)


def set_cell_text(cell, text, *, bold=False, color=None, size=10, fill=None):
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    set_run(run, size=size, bold=bold, color=color or DARK)
    if fill:
        shade_cell(cell, fill)


def academic_footer_docx(doc) -> None:
    add_para(doc, "", space_after=4)
    add_para(doc, ACADEMIC, size=9, color=MUTED, italic=True)


# ---------------------------------------------------------------------------
# 1. Correo PDF
# ---------------------------------------------------------------------------
def build_correo_pdf(path: Path) -> None:
    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=1.8 * cm,
        rightMargin=1.8 * cm,
        topMargin=1.5 * cm,
        bottomMargin=1.5 * cm,
    )
    styles = getSampleStyleSheet()
    brand = ParagraphStyle(
        "brand",
        parent=styles["Normal"],
        fontName="Helvetica-Bold",
        fontSize=9,
        textColor=colors.HexColor("#FFEC00"),
        spaceAfter=4,
    )
    title = ParagraphStyle(
        "title",
        parent=styles["Normal"],
        fontName="Helvetica-Bold",
        fontSize=14,
        textColor=colors.white,
        spaceAfter=0,
    )
    meta_label = ParagraphStyle(
        "meta_label",
        parent=styles["Normal"],
        fontName="Helvetica-Bold",
        fontSize=9,
        textColor=colors.HexColor("#4B3FAA"),
    )
    meta_val = ParagraphStyle(
        "meta_val",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=9,
        textColor=colors.HexColor("#211D40"),
    )
    body = ParagraphStyle(
        "body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#211D40"),
        spaceAfter=8,
        alignment=TA_JUSTIFY,
    )
    bullet = ParagraphStyle(
        "bullet",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=13,
        textColor=colors.HexColor("#211D40"),
        leftIndent=12,
        spaceAfter=2,
    )
    fiction = ParagraphStyle(
        "fiction",
        parent=styles["Normal"],
        fontName="Helvetica-Oblique",
        fontSize=8,
        textColor=colors.HexColor("#5A5A72"),
        alignment=TA_CENTER,
    )

    story = []

    header_data = [
        [Paragraph("OUTLOOK · MENSAJE DE PRÁCTICA", brand)],
        [Paragraph("SOLICITUD DE ANÁLISIS | Proyecto Horizonte", title)],
    ]
    header = Table(header_data, colWidths=[17 * cm])
    header.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#4B3FAA")),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("LEFTPADDING", (0, 0), (-1, -1), 12),
                ("RIGHTPADDING", (0, 0), (-1, -1), 12),
            ]
        )
    )
    bar = Table([[""]], colWidths=[17 * cm], rowHeights=[6])
    bar.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#211D40"))]))
    story.append(bar)
    story.append(header)
    story.append(Spacer(1, 6))

    meta_rows = [
        [
            Paragraph("De:", meta_label),
            Paragraph("Laura Méndez &lt;laura.mendez@ejemplo-academia.local&gt;", meta_val),
        ],
        [
            Paragraph("Para:", meta_label),
            Paragraph("Equipo de análisis de proyectos", meta_val),
        ],
        [
            Paragraph("Asunto:", meta_label),
            Paragraph("<b>SOLICITUD DE ANÁLISIS | Proyecto Horizonte</b>", meta_val),
        ],
        [
            Paragraph("Fecha:", meta_label),
            Paragraph("04/08/2026 09:15 (ficticia)", meta_val),
        ],
        [
            Paragraph("Área:", meta_label),
            Paragraph("Dirección de Infraestructura", meta_val),
        ],
    ]
    meta = Table(meta_rows, colWidths=[2.2 * cm, 14.8 * cm])
    meta.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F4F3F9")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#D0CCE0")),
            ]
        )
    )
    story.append(meta)
    story.append(Spacer(1, 14))

    story.append(Paragraph("Buenos días,", body))
    story.append(
        Paragraph(
            "Solicitamos realizar la revisión inicial del Proyecto Horizonte, "
            "iniciativa orientada a la modernización de infraestructura energética urbana.",
            body,
        )
    )
    story.append(
        Paragraph("Para este análisis se adjuntan los siguientes documentos:", body)
    )
    for item in [
        "Documento preliminar de alcance.",
        "Presupuesto y cronograma.",
        "Transcripción de la reunión inicial.",
        "Registro preliminar de riesgos.",
        "Comentarios de los interesados.",
    ]:
        story.append(Paragraph(f"• {item}", bullet))

    story.append(Spacer(1, 6))
    story.append(Paragraph("El comité requiere recibir:", body))
    for i, item in enumerate(
        [
            "Resumen ejecutivo.",
            "Alcance consolidado.",
            "Hallazgos presupuestales.",
            "Riesgos prioritarios.",
            "Compromisos identificados.",
            "Información faltante.",
            "Decisiones requeridas.",
            "Presentación ejecutiva.",
        ],
        1,
    ):
        story.append(Paragraph(f"{i}. {item}", bullet))

    story.append(Spacer(1, 8))
    story.append(
        Paragraph(
            "Por favor, no asuma que la iniciativa está aprobada. Diferencie "
            "claramente la información confirmada, las propuestas y los "
            "elementos pendientes de validación.",
            body,
        )
    )
    story.append(Spacer(1, 10))
    story.append(Paragraph("Cordialmente,", body))
    story.append(
        Paragraph(
            "<b>Laura Méndez</b><br/>Dirección de Infraestructura",
            ParagraphStyle(
                "sign",
                parent=body,
                leading=13,
                spaceAfter=12,
            ),
        )
    )

    attach = Table(
        [
            [
                Paragraph(
                    "<b>Adjuntos referenciados (kit de práctica):</b><br/>"
                    "02_Alcance · 03_Presupuesto_y_Cronograma · 04_Transcripcion · "
                    "05_Registro_Inicial_Riesgos · 06_Comentarios_Interesados",
                    ParagraphStyle(
                        "att",
                        parent=styles["Normal"],
                        fontName="Helvetica",
                        fontSize=8,
                        textColor=colors.HexColor("#211D40"),
                        leading=11,
                    ),
                )
            ]
        ],
        colWidths=[17 * cm],
    )
    attach.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F4F3F9")),
                ("BOX", (0, 0), (-1, -1), 1, colors.HexColor("#4B3FAA")),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("RIGHTPADDING", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    story.append(attach)
    story.append(Spacer(1, 16))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#211D40")))
    story.append(Spacer(1, 6))
    story.append(
        Paragraph(
            "Correo ficticio para fines académicos. · " + ACADEMIC,
            fiction,
        )
    )

    doc.build(story)


# ---------------------------------------------------------------------------
# 2. Alcance DOCX
# ---------------------------------------------------------------------------
def build_alcance_docx(path: Path) -> None:
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

    add_banner(
        doc,
        "Documento preliminar de alcance · Proyecto Horizonte",
        "Código PRY-HZ-2026 · Versión 0.3 · " + ACADEMIC,
    )

    meta = [
        ("Nombre del proyecto", "Proyecto Horizonte · Modernización de infraestructura energética urbana"),
        ("Código", "PRY-HZ-2026"),
        ("Patrocinador", "Dirección de Infraestructura"),
        ("Líder de proyecto", "Carlos Ramírez"),
        ("Presupuesto preliminar", "2.850.000.000 COP"),
        ("Fecha de inicio tentativa", "15/09/2026"),
        ("Fecha de fin tentativa", "30/06/2027"),
        ("Estado del documento", "Borrador para análisis · No aprobado"),
        ("Versión del plan de comunicación", "No especificado"),
        ("Gerente de contrato designado", "No especificado"),
    ]
    table = doc.add_table(rows=len(meta), cols=2)
    table.style = "Table Grid"
    for i, (k, v) in enumerate(meta):
        row = table.rows[i]
        set_cell_text(row.cells[0], k, bold=True, color=WHITE, size=9, fill=PURPLE)
        set_cell_text(row.cells[1], v, size=9, fill=LIGHT if i % 2 == 0 else WHITE)
        row.cells[0].width = Cm(5.5)
        row.cells[1].width = Cm(11.2)

    add_heading_styled(doc, "1. Necesidad")
    add_para(
        doc,
        "La zona urbana de referencia presenta crecimiento de demanda y "
        "envejecimiento de activos de distribución. Se requiere modernizar "
        "infraestructura crítica para mejorar confiabilidad, reducir interrupciones "
        "y preparar la red para mayor densidad de carga. La iniciativa se encuentra "
        "en fase de análisis preliminar; no constituye autorización de ejecución.",
        size=10,
    )

    add_heading_styled(doc, "2. Objetivo general")
    add_para(
        doc,
        "Modernizar la infraestructura energética urbana asociada al Proyecto "
        "Horizonte para elevar la confiabilidad del servicio, bajo criterios de "
        "seguridad, sostenibilidad y control presupuestal, sujeto a aprobación "
        "formal del comité directivo.",
        size=10,
    )

    add_heading_styled(doc, "3. Objetivos específicos")
    for t in [
        "Diagnosticar el estado técnico de equipos y obra civil relevantes.",
        "Actualizar equipos prioritarios según el plan de modernización.",
        "Definir controles de seguridad, HSE y gestión ambiental aplicables.",
        "Establecer un marco de seguimiento de cronograma y presupuesto.",
        "Preparar insumos para decisión del comité (sin asumir aprobación).",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "4. Alcance incluido")
    for t in [
        "Diagnóstico técnico preliminar de la instalación urbana de referencia.",
        "Ingeniería de detalle para las partidas priorizadas en el presupuesto.",
        "Suministro e instalación de equipos de potencia y control listados.",
        "Obra civil asociada a montaje y adecuaciones menores.",
        "Plan preliminar de riesgos, HSE y comunicación con interesados.",
        "Informes de avance y paquete de decisión para comité.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "5. Exclusiones")
    for t in [
        "Ampliar la red a municipios fuera del perímetro urbano definido.",
        "Integración con sistemas SCADA de terceros no autorizados.",
        "Compra de terrenos adicionales (No especificado si se requerirá).",
        "Operación comercial post-puesta en servicio (fuera de este alcance).",
        "Campañas publicitarias masivas (No especificado).",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "6. Entregables")
    for t in [
        "Informe de diagnóstico técnico.",
        "Paquete de ingeniería de detalle (revisión A).",
        "Plan de ejecución y cronograma actualizado.",
        "Control presupuestal consolidado.",
        "Matriz de riesgos priorizada.",
        "Informe de gestión de interesados.",
        "Acta de cierre de fase de análisis (si el comité lo autoriza).",
        "Presentación ejecutiva para comité.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "7. Hitos")
    hitos = [
        ("H1 · Kick-off interno", "15/09/2026", "Tentativo"),
        ("H2 · Diagnóstico cerrado", "30/10/2026", "Tentativo"),
        ("H3 · Ingeniería revisión A", "15/12/2026", "Tentativo"),
        ("H4 · Inicio de obra / montaje", "01/02/2027", "Tentativo"),
        ("H5 · Pruebas y energización", "15/05/2027", "Tentativo"),
        ("H6 · Cierre de proyecto", "30/06/2027", "Tentativo"),
        ("H7 · Auditoría externa de cierre", "No especificado", "No especificado"),
    ]
    ht = doc.add_table(rows=1 + len(hitos), cols=3)
    ht.style = "Table Grid"
    for j, h in enumerate(["Hito", "Fecha", "Estado"]):
        set_cell_text(ht.rows[0].cells[j], h, bold=True, color=WHITE, size=9, fill=DARK)
    for i, (n, f, e) in enumerate(hitos, 1):
        set_cell_text(ht.rows[i].cells[0], n, size=9)
        set_cell_text(ht.rows[i].cells[1], f, size=9)
        set_cell_text(ht.rows[i].cells[2], e, size=9, fill=LIGHT)

    add_heading_styled(doc, "8. Restricciones")
    for t in [
        "Presupuesto preliminar de referencia: 2.850.000.000 COP (sujeto a revisión).",
        "Ventana de intervención en horario diurno preferente; trabajos nocturnos requieren permiso adicional.",
        "No usar datos reales de clientes ni credenciales productivas en la práctica.",
        "Toda decisión de ejecución requiere aprobación humana del comité.",
        "Límite de personal contratista en sitio simultáneo: No especificado.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "9. Supuestos")
    for t in [
        "Se mantendrá acceso a la instalación durante el diagnóstico.",
        "Los precios de equipos se mantendrán dentro de ±8% de la estimación.",
        "La comunidad recibirá información previa a intervenciones mayores.",
        "Disponibilidad de ventanas de corte: No especificado (pendiente de Operaciones).",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "10. Dependencias")
    for t in [
        "Aprobación de presupuesto por Finanzas y comité directivo.",
        "Permisos ambientales / municipales aplicables.",
        "Disponibilidad de equipos de protección personal y protocolos HSE.",
        "Definición del contratista principal: No especificado.",
        "Liberación de planos as-built por Operaciones.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "11. Interesados")
    interesados = [
        ("Dirección de Infraestructura", "Patrocinio", "Laura Méndez"),
        ("Liderazgo de proyecto", "Ejecución", "Carlos Ramírez"),
        ("Operaciones", "Usuario técnico", "Andrés Quintero"),
        ("Finanzas", "Control presupuestal", "María Fernanda López"),
        ("HSE", "Seguridad y ambiente", "Patricia Gómez"),
        ("Comunidad zona urbana", "Afectados / vecinos", "Comité barrial (representante No especificado)"),
        ("Regulación / cumplimiento", "Cumplimiento", "No especificado"),
    ]
    it = doc.add_table(rows=1 + len(interesados), cols=3)
    it.style = "Table Grid"
    for j, h in enumerate(["Interesado", "Rol", "Contacto"]):
        set_cell_text(it.rows[0].cells[j], h, bold=True, color=WHITE, size=9, fill=PURPLE)
    for i, rowv in enumerate(interesados, 1):
        for j, val in enumerate(rowv):
            set_cell_text(it.rows[i].cells[j], val, size=9)

    add_heading_styled(doc, "12. Requisitos técnicos y de seguridad")
    for t in [
        "Cumplir normas técnicas internas de montaje y puesta en servicio.",
        "Bloqueo/etiquetado (LOTO) en intervenciones energizadas.",
        "Personal certificado para trabajos en altura y eléctricos.",
        "Especificación de marca preferente de relés: No especificado.",
        "Nivel de redundancia N-1 objetivo: No especificado.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "13. Requisitos ambientales")
    for t in [
        "Gestión de residuos de obra y aceites dieléctricos según procedimiento interno.",
        "Control de ruido y polvo en franja urbana.",
        "Estudio de impacto ambiental detallado: No especificado (verificar umbral legal).",
        "Plan de compensación arborícola: No especificado.",
    ]:
        add_bullet(doc, t)

    add_heading_styled(doc, "14. Comunicación")
    add_para(
        doc,
        "Se prevé comunicación periódica con el patrocinador y reporte semanal "
        "al líder de proyecto. Canal preferente: correo institucional y reuniones "
        "Teams. Frecuencia de boletín a la comunidad: No especificado. "
        "Protocolo de crisis mediática: No especificado.",
        size=10,
    )

    add_heading_styled(doc, "15. Criterios de aceptación")
    for t in [
        "Diagnóstico técnico revisado y firmado por Ingeniería.",
        "Desviación presupuestal explicada y validada por Finanzas.",
        "Riesgos prioritarios con controles o planes de acción asignados.",
        "Presentación de comité con fuentes trazables y datos no inventados.",
        "Indicador de disponibilidad post-modernización objetivo: No especificado.",
    ]:
        add_bullet(doc, t)

    academic_footer_docx(doc)
    doc.save(path)


# ---------------------------------------------------------------------------
# 3. Presupuesto y cronograma XLSX
# ---------------------------------------------------------------------------
def build_presupuesto_xlsx(path: Path) -> None:
    wb = Workbook()
    thin = Border(
        left=Side(style="thin", color="D0CCE0"),
        right=Side(style="thin", color="D0CCE0"),
        top=Side(style="thin", color="D0CCE0"),
        bottom=Side(style="thin", color="D0CCE0"),
    )
    fill_h = PatternFill("solid", fgColor=PURPLE)
    fill_d = PatternFill("solid", fgColor=DARK)
    fill_y = PatternFill("solid", fgColor=YELLOW)
    fill_l = PatternFill("solid", fgColor=LIGHT)
    fill_warn = PatternFill("solid", fgColor="FFE0E0")
    font_h = Font(name="Calibri", bold=True, color=WHITE, size=11)
    font_t = Font(name="Calibri", bold=True, color=YELLOW, size=13)
    font_n = Font(name="Calibri", color=DARK, size=10)
    font_b = Font(name="Calibri", bold=True, color=DARK, size=10)
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def hdr(ws, headers, row=3):
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=row, column=c, value=h)
            cell.fill = fill_h
            cell.font = font_h
            cell.alignment = center
            cell.border = thin

    def widths(ws, wsizes):
        for i, w in enumerate(wsizes, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    # --- Presupuesto ---
    # Word dice 2.850.000.000; Excel totaliza ~2.920.000.000 (inconsistencia intencional)
    # Una celda vacía y un valor atípico (Contingencia muy baja / partida errónea)
    ws = wb.active
    ws.title = "Presupuesto"
    ws["A1"] = "MCP-365-S2 · Presupuesto preliminar · Proyecto Horizonte (PRY-HZ-2026)"
    ws["A1"].font = font_t
    ws["A1"].fill = fill_d
    ws.merge_cells("A1:G1")
    ws.row_dimensions[1].height = 26
    ws["A2"] = (
        "Referencia en alcance Word: 2.850.000.000 COP · "
        "Este libro puede diferir levemente (práctica de contraste). · " + ACADEMIC
    )
    ws["A2"].font = Font(name="Calibri", italic=True, color=MUTED, size=9)
    ws.merge_cells("A2:G2")

    headers = [
        "Código",
        "Categoría",
        "Partida",
        "Monto (COP)",
        "Responsable",
        "Tipo de cifra",
        "Observación",
    ]
    hdr(ws, headers)

    rows = [
        ("HZ-ING-01", "Ingeniería", "Diagnóstico técnico e inspección", 185_000_000, "Carlos Ramírez", "Estimado", ""),
        ("HZ-ING-02", "Ingeniería", "Ingeniería de detalle revisión A", 320_000_000, "Ingeniería", "Estimado", ""),
        ("HZ-EQ-01", "Equipos", "Transformadores y celdas", 980_000_000, "Compras", "Estimado", ""),
        ("HZ-EQ-02", "Equipos", "Protecciones, control y comunicaciones", 410_000_000, "Compras", "Estimado", ""),
        ("HZ-OB-01", "Obra civil", "Adecuaciones de obra civil y montaje", 520_000_000, "Construcción", "Estimado", ""),
        ("HZ-OB-02", "Obra civil", "Obra civil complementaria periurbana", 45_000_000, "Construcción", "Atípico", "Valor atípico bajo vs. alcance verbal — revisar"),
        ("HZ-SEG-01", "Seguridad", "HSE, EPP y señalización", 95_000_000, "Patricia Gómez", "Estimado", ""),
        ("HZ-SEG-02", "Seguridad", "Estudios de arco eléctrico", None, "HSE", "No especificado", "Celda vacía intencional — No especificado"),
        ("HZ-CON-01", "Contingencia", "Reserva de contingencia (5% ref.)", 145_000_000, "Finanzas", "Estimado", "Por debajo de práctica usual 8-10%"),
        ("HZ-GES-01", "Gestión", "Gestión de proyecto y PMO", 120_000_000, "Carlos Ramírez", "Estimado", ""),
        ("HZ-GES-02", "Gestión", "Comunicación con interesados", 55_000_000, "Laura Méndez", "Estimado", ""),
        ("HZ-GES-03", "Gestión", "Auditoría y aseguramiento de calidad", 45_000_000, "No especificado", "Estimado", ""),
    ]

    total = 0
    for i, row in enumerate(rows, 4):
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=i, column=c, value=val if val is not None else "")
            cell.font = font_n
            cell.alignment = left if c != 4 else Alignment(horizontal="right", vertical="center")
            cell.border = thin
            cell.fill = fill_l if i % 2 == 0 else PatternFill()
            if c == 4 and isinstance(val, (int, float)):
                cell.number_format = '#,##0'
                total += val
            if row[6] and ("vacía" in str(row[6]).lower() or "atípico" in str(row[6]).lower()):
                if c == 4:
                    cell.fill = fill_warn

    r_tot = 4 + len(rows)
    ws.cell(row=r_tot, column=3, value="TOTAL (suma de montos informados)").font = font_b
    cell_t = ws.cell(row=r_tot, column=4, value=total)
    cell_t.font = Font(name="Calibri", bold=True, color=DARK, size=11)
    cell_t.fill = fill_y
    cell_t.number_format = '#,##0'
    cell_t.border = thin
    ws.cell(row=r_tot + 1, column=1, value=(
        f"Total libro ≈ {total:,} COP. Diferencia vs. Word (2.850.000.000): "
        f"práctica de detección de inconsistencias. Partida HZ-SEG-02 sin monto."
    )).font = Font(name="Calibri", italic=True, color=MUTED, size=9)
    ws.merge_cells(start_row=r_tot + 1, start_column=1, end_row=r_tot + 1, end_column=7)
    widths(ws, [12, 14, 42, 16, 16, 14, 40])

    # --- Cronograma ---
    # Word H4 = 01/02/2027; Excel pone 08/02/2027 (1 semana de diferencia intencional)
    ws2 = wb.create_sheet("Cronograma")
    ws2["A1"] = "MCP-365-S2 · Cronograma preliminar · Proyecto Horizonte"
    ws2["A1"].font = font_t
    ws2["A1"].fill = fill_d
    ws2.merge_cells("A1:H1")
    ws2.row_dimensions[1].height = 26
    ws2["A2"] = (
        "Fechas tentativas. Ruta crítica marcada. "
        "Hito de inicio de obra difiere ~1 semana vs. documento de alcance. · " + ACADEMIC
    )
    ws2["A2"].font = Font(name="Calibri", italic=True, color=MUTED, size=9)
    ws2.merge_cells("A2:H2")

    ch = [
        "ID",
        "Actividad",
        "Inicio",
        "Fin",
        "Duración (días)",
        "Predecesora",
        "Ruta crítica",
        "Responsable",
    ]
    hdr(ws2, ch)
    acts = [
        ("A1", "Kick-off y alineación de alcance", "2026-09-15", "2026-09-22", 6, "—", "Sí", "Carlos Ramírez"),
        ("A2", "Diagnóstico técnico en sitio", "2026-09-23", "2026-10-30", 28, "A1", "Sí", "Ingeniería"),
        ("A3", "Ingeniería de detalle (Rev. A)", "2026-11-02", "2026-12-15", 32, "A2", "Sí", "Ingeniería"),
        ("A4", "Proceso de compras de equipos", "2026-11-10", "2027-01-20", 50, "A2", "Sí", "Compras"),
        ("A5", "Gestión de permisos y HSE", "2026-10-01", "2027-01-15", 75, "A1", "No", "Patricia Gómez"),
        ("A6", "Inicio de obra / montaje", "2027-02-08", "2027-04-30", 60, "A3,A4", "Sí", "Construcción"),
        ("A7", "Pruebas, energización y ajustes", "2027-05-03", "2027-05-28", 20, "A6", "Sí", "Operaciones"),
        ("A8", "Comunicación con comunidad", "2026-09-20", "2027-06-15", 190, "A1", "No", "Laura Méndez"),
        ("A9", "Cierre documental y lecciones", "2027-06-01", "2027-06-30", 22, "A7", "No", "Carlos Ramírez"),
        ("A10", "Auditoría externa de cierre", "No especificado", "No especificado", "—", "A9", "No", "No especificado"),
    ]
    for i, row in enumerate(acts, 4):
        for c, val in enumerate(row, 1):
            cell = ws2.cell(row=i, column=c, value=val)
            cell.font = font_n
            cell.alignment = center if c in (1, 5, 6, 7) else left
            cell.border = thin
            if row[6] == "Sí":
                cell.fill = PatternFill("solid", fgColor="EDE9FF")
            if row[0] == "A6":
                cell.fill = fill_warn
    ws2["A15"] = (
        "Nota de contraste: en 02_Alcance el hito H4 (Inicio de obra/montaje) "
        "aparece como 01/02/2027; aquí A6 inicia 08/02/2027."
    )
    ws2["A15"].font = Font(name="Calibri", italic=True, color=MUTED, size=9)
    ws2.merge_cells("A15:H15")
    widths(ws2, [6, 38, 14, 14, 14, 12, 12, 16])

    # --- Notas ---
    ws3 = wb.create_sheet("Notas")
    ws3["A1"] = "Notas de uso · Proyecto Horizonte"
    ws3["A1"].font = font_t
    ws3["A1"].fill = fill_d
    ws3.merge_cells("A1:B1")
    notes = [
        ("Propósito", "Fuente de práctica para contrastar cifras y fechas con Word y transcripción."),
        ("Inconsistencias intencionales", "Total ≠ 2.850.000.000; A6 vs H4; celda vacía HZ-SEG-02; partida atípica HZ-OB-02."),
        ("Aprobación", "Ninguna cifra de este libro implica aprobación del proyecto."),
        ("Vacíos", "Usar «No especificado» donde no haya evidencia. No inventar."),
        ("Etiqueta", ACADEMIC),
    ]
    ws3["A3"] = "Tema"
    ws3["B3"] = "Detalle"
    ws3["A3"].fill = fill_h
    ws3["B3"].fill = fill_h
    ws3["A3"].font = font_h
    ws3["B3"].font = font_h
    for i, (a, b) in enumerate(notes, 4):
        ws3.cell(row=i, column=1, value=a).font = font_b
        ws3.cell(row=i, column=2, value=b).font = font_n
        ws3.cell(row=i, column=1).border = thin
        ws3.cell(row=i, column=2).border = thin
    widths(ws3, [28, 90])

    wb.save(path)


# ---------------------------------------------------------------------------
# 4. Transcripción DOCX
# ---------------------------------------------------------------------------
def build_transcripcion_docx(path: Path) -> None:
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

    add_banner(
        doc,
        "Transcripción · Reunión inicial Proyecto Horizonte",
        "Teams · 28/07/2026 10:00–11:15 · " + ACADEMIC,
    )
    add_para(
        doc,
        "Participantes: Laura Méndez (Infraestructura), Carlos Ramírez (Líder), "
        "Andrés Quintero (Operaciones), María Fernanda López (Finanzas), "
        "Patricia Gómez (HSE), Sofía Herrera (PMO · facilitadora).",
        size=9,
        color=MUTED,
    )
    add_para(
        doc,
        "Nota: transcripción condensada y ficticia. Mezcla decisiones, propuestas, "
        "opiniones, preguntas, preocupaciones y tareas. No equivale a un acta aprobada.",
        size=9,
        italic=True,
        color=MUTED,
    )

    turns = [
        ("Sofía Herrera", "Abrimos la reunión de alineación del Proyecto Horizonte. "
         "Objetivo: revisar necesidad, cifras preliminares y pendientes para el comité. "
         "Recuerdo: no hay aprobación de ejecución en esta sesión."),
        ("Laura Méndez", "DECISIÓN CONFIRMADA: solicitamos formalmente el análisis "
         "integrado con los cinco anexos. El patrocinio queda en Dirección de Infraestructura."),
        ("Carlos Ramírez", "PROPUESTA: arrancar kick-off el 15/09/2026. Esa fecha la "
         "tengo como tentativa hasta confirmar disponibilidad de Operaciones."),
        ("Andrés Quintero", "PREGUNTA: ¿ya está definido el número de ventanas de corte? "
         "Sin eso no puedo comprometer el cronograma de montaje."),
        ("Carlos Ramírez", "Respuesta: aún No especificado. Lo dejamos como dependencia."),
        ("María Fernanda López", "OPINIÓN / PROPUESTA: el presupuesto de 2.850.000.000 COP "
         "se ve justo. Sugiero una revisión presupuestal al alza hacia 3.050.000.000 COP "
         "para cubrir contingencia realista y estudios de arco. No es una decisión; es "
         "una propuesta para el comité."),
        ("Laura Méndez", "Anotado. No aprobamos ese monto hoy. Mantener el preliminar "
         "del alcance y marcar la propuesta de Finanzas como pendiente de validación."),
        ("Patricia Gómez", "PREOCUPACIÓN: en el registro de riesgos veo controles vacíos. "
         "No inventemos mitigaciones. Prefiero «No especificado» hasta visita de campo."),
        ("Andrés Quintero", "TAREA ASIGNADA (confirmada): Operaciones entregará planos "
         "as-built parciales a más tardar el 20/08/2026. Responsable: Andrés Quintero."),
        ("Sofía Herrera", "TAREA TENTATIVA: PMO preparará borrador de presentación "
         "para comité la semana del 10/08/2026, dueño sugerido Sofía Herrera — "
         "sujeto a carga de otros proyectos."),
        ("Carlos Ramírez", "Sobre el hito de inicio de obra: en el alcance puse 01/02/2027. "
         "En el Excel de cronograma estoy moviendo a 08/02/2027 por dependencia de compras. "
         "Hay que transparentar esa diferencia."),
        ("María Fernanda López", "PREGUNTA: ¿quién firma el umbral de desviación presupuestal? "
         "Contacto de regulación / cumplimiento sigue No especificado."),
        ("Laura Méndez", "Correcto. No inventar nombres. Llevar esa pregunta al comité."),
        ("Patricia Gómez", "PROPUESTA no aprobada: incluir un simulacro de emergencia "
         "comunitaria antes del montaje. Operaciones dice que puede ser excesivo en esta fase."),
        ("Andrés Quintero", "OPINIÓN: de acuerdo en diferir el simulacro. Prioricemos LOTO "
         "y señalización."),
        ("Sofía Herrera", "Compromiso confirmado de esta reunión: (1) análisis con anexos; "
         "(2) entrega de planos parciales 20/08/2026; (3) no asumir aprobación del proyecto; "
         "(4) documentar propuesta de techo 3.050.000.000 COP como no aprobada."),
        ("Carlos Ramírez", "Cierro con próximo paso: consolidar hallazgos para el correo "
         "de Laura y la plantilla de comité. Fecha de comité: tentativa 18/08/2026, "
         "no confirmada por agenda de dirección."),
        ("Sofía Herrera", "Gracias. Fin de la grabación. Material con fines académicos."),
    ]

    for speaker, text in turns:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(8)
        r1 = p.add_run(f"{speaker}: ")
        set_run(r1, size=10, bold=True, color=PURPLE)
        r2 = p.add_run(text)
        set_run(r2, size=10, color=DARK)

    add_heading_styled(doc, "Leyenda rápida para el análisis")
    for t in [
        "Confirmado: patrocinio, solicitud de análisis, tarea de planos 20/08/2026.",
        "Tentativo: kick-off 15/09/2026, comité 18/08/2026, borrador PPT semana 10/08.",
        "Propuesta no aprobada: techo presupuestal 3.050.000.000 COP; simulacro comunitario.",
        "Contradicción leve vs. alcance/presupuesto: revisión al alza y fecha de obra 08/02 vs 01/02.",
        "Vacíos: ventanas de corte, firmante de umbral, regulación — No especificado.",
    ]:
        add_bullet(doc, t)

    academic_footer_docx(doc)
    doc.save(path)


# ---------------------------------------------------------------------------
# 5. Riesgos XLSX
# ---------------------------------------------------------------------------
def build_riesgos_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Riesgos"
    thin = Border(
        left=Side(style="thin", color="D0CCE0"),
        right=Side(style="thin", color="D0CCE0"),
        top=Side(style="thin", color="D0CCE0"),
        bottom=Side(style="thin", color="D0CCE0"),
    )
    fill_h = PatternFill("solid", fgColor=PURPLE)
    fill_d = PatternFill("solid", fgColor=DARK)
    fill_l = PatternFill("solid", fgColor=LIGHT)
    font_h = Font(name="Calibri", bold=True, color=WHITE, size=10)
    font_t = Font(name="Calibri", bold=True, color=YELLOW, size=12)
    font_n = Font(name="Calibri", color=DARK, size=9)
    wrap = Alignment(wrap_text=True, vertical="center")

    ws["A1"] = "Registro inicial de riesgos · Proyecto Horizonte (PRY-HZ-2026)"
    ws["A1"].font = font_t
    ws["A1"].fill = fill_d
    ws.merge_cells("A1:J1")
    ws.row_dimensions[1].height = 24
    ws["A2"] = (
        "Valoración preliminar no definitiva. Controles vacíos = No especificado. · "
        + ACADEMIC
    )
    ws["A2"].font = Font(name="Calibri", italic=True, color=MUTED, size=9)
    ws.merge_cells("A2:J2")

    headers = [
        "ID",
        "Riesgo",
        "Categoría",
        "Fuente",
        "Causa",
        "Consecuencia",
        "Prob. preliminar",
        "Impacto preliminar",
        "Control existente",
        "Notas",
    ]
    for c, h in enumerate(headers, 1):
        cell = ws.cell(row=3, column=c, value=h)
        cell.fill = fill_h
        cell.font = font_h
        cell.alignment = Alignment(wrap_text=True, horizontal="center", vertical="center")
        cell.border = thin
    ws.row_dimensions[3].height = 30

    risks = [
        ("R-HZ-01", "Retraso en ventanas de corte no definidas", "Operativo",
         "Transcripción / Alcance", "Dependencia con Operaciones sin fecha",
         "Desplazamiento del montaje y ruta crítica", "Media", "Alto",
         "No especificado", "Explícito como vacío"),
        ("R-HZ-02", "Subestimación de contingencia presupuestal", "Financiero",
         "Presupuesto / Transcripción", "Reserva ~5% y propuesta de alza no aprobada",
         "Desvío de costo o recorte de alcance", "Media", "Alto",
         "Revisión preliminar Finanzas", "Contraste Word vs Excel"),
        ("R-HZ-03", "Incidentes HSE en zona urbana densa", "Seguridad",
         "Comentarios interesados / HSE", "Trabajos cerca de comunidad y tráfico",
         "Lesiones, paradas y reputación", "Media", "Muy alto",
         "No especificado", "Control pendiente de visita"),
        ("R-HZ-04", "Rechazo comunitario por ruido/polvo", "Social",
         "Comentarios interesados", "Comunicación a comunidad No especificada",
         "Bloqueos, quejas y retrasos de permiso", "Media", "Alto",
         "No especificado", ""),
        ("R-HZ-05", "Demora en suministro de transformadores", "Cronograma",
         "Cronograma Excel", "Lead time de compras en ruta crítica",
         "Atraso de hito de montaje", "Alta", "Alto",
         "Seguimiento Compras (preliminar)", ""),
        ("R-HZ-06", "Inconsistencia de fechas entre documentos", "Contractual",
         "Alcance vs Cronograma", "H4 01/02 vs A6 08/02",
         "Confusión en compromisos reportados al comité", "Alta", "Medio",
         "No especificado", "Para práctica de auditoría"),
        ("R-HZ-07", "Vacíos en requisitos ambientales formales", "Ambiental",
         "Alcance", "Estudio de impacto No especificado",
         "Hallazgo de cumplimiento o freno de obra", "Baja", "Alto",
         "No especificado", ""),
        ("R-HZ-08", "Falta de firmante de umbral presupuestal", "Financiero",
         "Transcripción", "Rol de regulación/cumplimiento vacío",
         "Decisiones sin autoridad clara", "Media", "Medio",
         "No especificado", ""),
        ("R-HZ-09", "Interferencia con operación en servicio", "Técnico",
         "Operaciones", "Modernización sobre activos energizados",
         "Interrupciones no planificadas", "Media", "Alto",
         "Procedimientos LOTO (referencia)", "Alcance parcial"),
        ("R-HZ-10", "Filtración de datos de práctica como reales", "Reputacional",
         "Diseño del caso académico", "Material didáctico mal contextualizado",
         "Malentendido externo", "Baja", "Medio",
         "Etiqueta académica en archivos", "Mitigación didáctica"),
    ]

    for i, row in enumerate(risks, 4):
        for c, val in enumerate(row, 1):
            cell = ws.cell(row=i, column=c, value=val)
            cell.font = font_n
            cell.alignment = wrap
            cell.border = thin
            cell.fill = fill_l if i % 2 == 0 else PatternFill()
            if c == 9 and val == "No especificado":
                cell.fill = PatternFill("solid", fgColor="FFF8CC")
        ws.row_dimensions[i].height = 42

    for i, w in enumerate([10, 36, 12, 18, 28, 28, 12, 12, 22, 22], 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    wb.save(path)


# ---------------------------------------------------------------------------
# 6. Comentarios interesados DOCX
# ---------------------------------------------------------------------------
def build_comentarios_docx(path: Path) -> None:
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

    add_banner(
        doc,
        "Comentarios de interesados · Proyecto Horizonte",
        "Recopilación preliminar · " + ACADEMIC,
    )
    add_para(
        doc,
        "Documento de escucha temprana. Las solicitudes no están aprobadas. "
        "Usar para identificar preocupaciones, requisitos y vacíos de información.",
        size=10,
        color=MUTED,
    )

    blocks = [
        (
            "1. Comunidad (zona urbana de referencia)",
            [
                "Preocupación: ruido nocturno y polvo durante obra civil.",
                "Solicitud: aviso con al menos 10 días de anticipación antes de cortes.",
                "Solicitud: canal único de PQRS con horario de atención publicado.",
                "Pregunta: ¿habrá reubicación temporal de puestos comerciales? Respuesta actual: No especificado.",
                "Opinión: apoyan modernización si se reduce duración de interrupciones históricas.",
            ],
        ),
        (
            "2. Operaciones (Andrés Quintero)",
            [
                "Preocupación: ventanas de corte aún no cuantificadas.",
                "Solicitud: planos as-built validados antes de ingeniería Rev. A.",
                "Solicitud: no solapar montaje con mantenimiento mayor programado en Q1-2027.",
                "Advertencia: partida de obra civil periurbana (HZ-OB-02) parece subestimada.",
                "Compromiso confirmado en reunión: entregar planos parciales el 20/08/2026.",
            ],
        ),
        (
            "3. Finanzas (María Fernanda López)",
            [
                "Preocupación: contingencia insuficiente frente a lead time de equipos.",
                "Propuesta (no aprobada): elevar techo a 3.050.000.000 COP para revisión de comité.",
                "Solicitud: marcar toda cifra como estimada hasta cotizaciones formales.",
                "Solicitud: definir firmante del umbral de desviación (hoy No especificado).",
                "Observación: diferencia entre total Excel y 2.850.000.000 del alcance debe explicarse.",
            ],
        ),
        (
            "4. HSE (Patricia Gómez)",
            [
                "Preocupación: controles existentes vacíos en varios riesgos del registro.",
                "Solicitud: visita de campo HSE antes de fijar probabilidad/impacto definitivos.",
                "Solicitud: plan de señalización peatonal en franja escolar cercana.",
                "Rechazo tentativo a simulacro comunitario en esta fase (propuesta diferida).",
                "Requisito: personal con certificación vigente para trabajos eléctricos y en altura.",
            ],
        ),
    ]

    for title, items in blocks:
        add_heading_styled(doc, title)
        for item in items:
            add_bullet(doc, item)

    add_heading_styled(doc, "5. Síntesis de pedidos recurrentes")
    for t in [
        "Transparentar fechas de corte y comunicación comunitaria.",
        "No presentar propuestas presupuestales como aprobadas.",
        "Completar controles HSE con evidencia de campo.",
        "Resolver vacíos etiquetados como No especificado antes del comité.",
    ]:
        add_bullet(doc, t)

    academic_footer_docx(doc)
    doc.save(path)


# ---------------------------------------------------------------------------
# 7. Plantilla PPTX
# ---------------------------------------------------------------------------
def build_plantilla_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_spec = [
        "1. Contexto y necesidad",
        "2. Objetivo y alcance",
        "3. Entregables e hitos",
        "4. Presupuesto",
        "5. Cronograma",
        "6. Riesgos prioritarios",
        "7. Decisiones requeridas",
        "8. Próximos pasos",
    ]

    blank = prs.slide_layouts[6]  # blank

    for title in slides_spec:
        slide = prs.slides.add_slide(blank)

        # Top bar
        bar = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.35),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PptRGB(0x21, 0x1D, 0x40)
        bar.line.fill.background()

        # Accent strip
        accent = slide.shapes.add_shape(
            1,
            Inches(0),
            Inches(0.35),
            Inches(0.18),
            Inches(7.15),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = PptRGB(0x4B, 0x3F, 0xAA)
        accent.line.fill.background()

        # Title
        tf = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.8))
        p = tf.text_frame.paragraphs[0]
        p.text = title
        p.font.size = PptPt(32)
        p.font.bold = True
        p.font.color.rgb = PptRGB(0x21, 0x1D, 0x40)
        p.font.name = "Calibri"

        # Subtitle placeholder
        sf = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(0.5))
        sp = sf.text_frame.paragraphs[0]
        sp.text = "Pendiente de validación"
        sp.font.size = PptPt(18)
        sp.font.italic = True
        sp.font.color.rgb = PptRGB(0x4B, 0x3F, 0xAA)
        sp.font.name = "Calibri"

        # Content placeholders
        body = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(12), Inches(3.5))
        bf = body.text_frame
        bf.word_wrap = True
        lines = [
            "• Mensaje principal: Pendiente de validación",
            "• Punto de evidencia: Pendiente de validación",
            "• Dato de soporte: Pendiente de validación",
            "• Información pendiente: Pendiente de validación",
        ]
        for i, line in enumerate(lines):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = line
            para.font.size = PptPt(16)
            para.font.color.rgb = PptRGB(0x5A, 0x5A, 0x72)
            para.font.name = "Calibri"
            para.space_after = PptPt(10)

        # Footer
        foot = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12), Inches(0.35))
        fp = foot.text_frame.paragraphs[0]
        fp.text = (
            "Proyecto Horizonte · PRY-HZ-2026 · Plantilla de comité · "
            + ACADEMIC
        )
        fp.font.size = PptPt(10)
        fp.font.color.rgb = PptRGB(0x5A, 0x5A, 0x72)
        fp.font.name = "Calibri"
        fp.alignment = PP_ALIGN.LEFT

        # Yellow corner mark
        mark = slide.shapes.add_shape(
            1,
            Inches(12.6),
            Inches(7.15),
            Inches(0.55),
            Inches(0.18),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = PptRGB(0xFF, 0xEC, 0x00)
        mark.line.fill.background()

    prs.save(str(path))


# ---------------------------------------------------------------------------
# 8. Guía validación PDF
# ---------------------------------------------------------------------------
def build_guia_pdf(path: Path) -> None:
    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=1.8 * cm,
        rightMargin=1.8 * cm,
        topMargin=1.6 * cm,
        bottomMargin=1.6 * cm,
    )
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle(
        "h1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=16,
        textColor=colors.HexColor("#211D40"),
        spaceAfter=10,
    )
    h2 = ParagraphStyle(
        "h2",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=12,
        textColor=colors.HexColor("#4B3FAA"),
        spaceBefore=10,
        spaceAfter=6,
    )
    body = ParagraphStyle(
        "body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=13,
        textColor=colors.HexColor("#211D40"),
        alignment=TA_JUSTIFY,
        spaceAfter=6,
    )
    bullet = ParagraphStyle(
        "b",
        parent=body,
        leftIndent=10,
        spaceAfter=3,
    )
    small = ParagraphStyle(
        "small",
        parent=body,
        fontSize=8,
        textColor=colors.HexColor("#5A5A72"),
        alignment=TA_CENTER,
    )

    story = []
    head = Table(
        [[Paragraph("MCP-365-S2 · Guía de validación de resultados · Proyecto Horizonte", 
                    ParagraphStyle("ht", parent=h1, textColor=colors.HexColor("#FFEC00"), fontSize=13))]],
        colWidths=[17 * cm],
    )
    head.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#211D40")),
                ("LEFTPADDING", (0, 0), (-1, -1), 12),
                ("TOPPADDING", (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ]
        )
    )
    story.append(head)
    story.append(Spacer(1, 10))
    story.append(
        Paragraph(
            "Cómo verificar las salidas de Copilot frente a las fuentes del caso. "
            "Todo el material es ficticio y de uso académico.",
            body,
        )
    )

    story.append(Paragraph("1. Principio de trazabilidad", h2))
    story.append(
        Paragraph(
            "Cada afirmación relevante debe poder señalarse en un archivo fuente "
            "(correo, alcance, presupuesto/cronograma, transcripción, riesgos o "
            "comentarios). Si no hay evidencia, no se completa con invención.",
            body,
        )
    )

    story.append(Paragraph("2. Diferenciar tipos de información", h2))
    for t in [
        "<b>Explícito:</b> aparece literalmente en una fuente (cifra, fecha, nombre, frase).",
        "<b>Inferido:</b> se deriva de combinar fuentes; debe etiquetarse como inferencia y validarse.",
        "<b>No especificado:</b> la fuente no lo trae o deja el campo vacío; escribir exactamente «No especificado».",
        "<b>Propuesta / tentativo:</b> sugerido en reunión o comentarios, sin aprobación.",
        "<b>Confirmado:</b> decisión o tarea acordada con evidencia en transcripción u otro documento.",
    ]:
        story.append(Paragraph(f"• {t}", bullet))

    story.append(Paragraph("3. Cómo contrastar salidas de Copilot", h2))
    for t in [
        "Abrir el archivo fuente citado y localizar la evidencia (sección, celda o turno).",
        "Comparar montos: Word (2.850.000.000) vs total Excel (puede diferir a propósito).",
        "Comparar fechas de hitos (p. ej. inicio de obra 01/02 vs 08/02).",
        "Separar la propuesta de techo 3.050.000.000 COP (no aprobada) del preliminar.",
        "Revisar que los controles de riesgo vacíos sigan como No especificado.",
        "Comprobar que el correo no se interprete como aprobación del proyecto.",
    ]:
        story.append(Paragraph(f"• {t}", bullet))

    story.append(Paragraph("4. Lista de comprobación", h2))
    checks = [
        "Las cifras citadas coinciden con al menos una fuente o se marcan como inconsistentes.",
        "Las fechas coinciden o se documenta la discrepancia.",
        "Los riesgos tienen fuente y no se inventan controles.",
        "Las decisiones están diferenciadas de propuestas y opiniones.",
        "Los datos faltantes están marcados como No especificado.",
        "No se afirma que el proyecto esté aprobado.",
        "La presentación usa «Pendiente de validación» donde falte evidencia.",
        "Una persona revisó el entregable antes de enviarlo al comité (simulado).",
    ]
    for t in checks:
        story.append(Paragraph(f"☐ {t}", bullet))

    story.append(Paragraph("5. Tratamiento de inconsistencias halladas", h2))
    story.append(
        Paragraph(
            "Registrar tema, archivos comparados, evidencia, tipo de inconsistencia, "
            "impacto potencial y pregunta de validación. No decidir unilateralmente "
            "cuál documento «gana»: elevar al responsable humano.",
            body,
        )
    )

    story.append(Spacer(1, 14))
    story.append(
        Paragraph(ACADEMIC + " · Proyecto Horizonte · PRY-HZ-2026", small)
    )

    doc.build(story)


# ---------------------------------------------------------------------------
# ZIP + main
# ---------------------------------------------------------------------------
def build_zip(out_dir: Path, zip_name: str = "MCP365_S2_Kit_Proyecto_Horizonte.zip") -> Path:
    zip_path = out_dir / zip_name
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name in FILENAMES:
            fp = out_dir / name
            if not fp.exists():
                raise FileNotFoundError(f"Falta archivo para ZIP: {fp}")
            zf.write(fp, arcname=name)
    return zip_path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    builders = [
        (FILENAMES[0], build_correo_pdf),
        (FILENAMES[1], build_alcance_docx),
        (FILENAMES[2], build_presupuesto_xlsx),
        (FILENAMES[3], build_transcripcion_docx),
        (FILENAMES[4], build_riesgos_xlsx),
        (FILENAMES[5], build_comentarios_docx),
        (FILENAMES[6], build_plantilla_pptx),
        (FILENAMES[7], build_guia_pdf),
    ]

    print(f"Salida: {OUT_DIR}")
    for name, fn in builders:
        target = OUT_DIR / name
        print(f"  Generando {name} …")
        fn(target)

    zip_path = build_zip(OUT_DIR)
    print(f"  ZIP: {zip_path.name}")

    print("\n=== Verificación ===")
    ok = True
    for name in FILENAMES:
        fp = OUT_DIR / name
        size = fp.stat().st_size if fp.exists() else 0
        status = "OK" if fp.exists() and size > 0 else "FAIL"
        if status == "FAIL":
            ok = False
        print(f"  [{status}] {name:52s} {size:>10,} bytes")
    zsize = zip_path.stat().st_size
    print(f"  [{'OK' if zsize > 0 else 'FAIL'}] {zip_path.name:52s} {zsize:>10,} bytes")
    if not ok:
        raise SystemExit(1)
    print("\nListo.")


if __name__ == "__main__":
    main()
